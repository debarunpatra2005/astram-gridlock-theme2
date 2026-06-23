"""
Generate an EDITABLE PowerPoint deck for the ASTRAM Theme 2 submission.

Same 8 slides as the PDF, but as .pptx so you can tweak text/branding.
Writes prototype/ASTRAM_Theme2_Pitch_Deck.pptx (16:9).

Run:  prototype/.venv/bin/python make_pptx.py
"""

import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
ART = os.path.join(HERE, "artifacts")
FIG = os.path.join(HERE, "figures")
OUT = os.path.join(HERE, "ASTRAM_Theme2_Pitch_Deck.pptx")

NAVY = RGBColor(0x0B, 0x25, 0x45)
BLUE = RGBColor(0x25, 0x63, 0xEB)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREY = RGBColor(0x47, 0x55, 0x69)
DARK = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCB, 0xD5, 0xE1)
CARD = RGBColor(0xF1, 0xF5, 0xF9)

with open(os.path.join(ART, "metrics.json")) as f:
    M = json.load(f)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=False):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                             Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after)
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = "Calibri"
    return tb


def header(s, kicker, title):
    rect(s, 1.0, 0.95, 0.55, 0.10, AMBER)
    text(s, 1.0, 0.45, 10, 0.5, [[(kicker, 15, BLUE, True)]])
    text(s, 1.0, 1.15, 11.3, 1.0, [[(title, 30, NAVY, True)]])


def bullets(s, items, y=2.5, dy=0.92, size=17):
    for i, it in enumerate(items):
        yy = y + i * dy
        text(s, 1.0, yy, 0.4, 0.5, [[("▸", size, BLUE, True)]])
        text(s, 1.45, yy, 10.6, dy, [[(it, size, DARK, False)]])


def pic(s, path, x, y, w):
    s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


# ---- Slide 1: Title -------------------------------------------------------
s = slide()
rect(s, 0, 0, 13.333, 7.5, NAVY)
text(s, 0.5, 2.2, 12.3, 1.2, [[("ASTRAM Congestion Copilot", 44, WHITE, True)]],
     align=PP_ALIGN.CENTER)
text(s, 0.5, 3.4, 12.3, 0.7,
     [[("Forecasting event-driven congestion — and the plan to clear it", 20,
        LIGHT, False)]], align=PP_ALIGN.CENTER)
rect(s, 5.5, 4.25, 2.3, 0.06, AMBER)
text(s, 0.5, 4.6, 12.3, 0.5,
     [[("Gridlock Hackathon 2.0  ·  Round 2 (Prototype)  ·  Theme 2", 15, WHITE,
        False)]], align=PP_ALIGN.CENTER)
text(s, 0.5, 5.15, 12.3, 0.5,
     [[("Event-Driven Congestion (Planned & Unplanned)", 14, LIGHT, False)]],
     align=PP_ALIGN.CENTER)

# ---- Slide 2: Problem -----------------------------------------------------
s = slide()
header(s, "THE PROBLEM", "Events break traffic faster than crews can react")
bullets(s, [
    "Rallies, festivals, sports, construction & sudden gatherings cause localized breakdowns.",
    "Today: event impact isn't quantified in advance.",
    "Resource deployment is experience-driven, not data-driven.",
    "No post-event learning system to improve next time.",
], y=2.4, dy=0.78)
text(s, 1.0, 5.9, 11.5, 1.0,
     [[("Theme 2 asks two things — forecast impact AND recommend manpower, "
        "barricading & diversion. Most stop at prediction. We do both.", 16,
        NAVY, True)]])

# ---- Slide 3: Solution ----------------------------------------------------
s = slide()
header(s, "OUR SOLUTION", "A working forecast-to-dispatch prototype")
labels = [("FORECAST", "clearance time as a P10 / P50 / P90 interval, not a fragile guess."),
          ("RECOMMEND", "response team, crew count, equipment, barricading & diversion."),
          ("DISPATCH", "shift-level pre-positioning plan by corridor × time window (CSV/JSON)."),
          ("PREVENT", "recurring waterlogging / pothole / construction hotspots.")]
for i, (k, v) in enumerate(labels):
    yy = 2.4 + i * 0.92
    text(s, 1.0, yy, 11.0, 0.8,
         [[("▸  ", 17, BLUE, True), (k + "  —  ", 17, NAVY, True),
           (v, 17, DARK, False)]])
text(s, 1.0, 6.4, 11.0, 0.5,
     [[("Live Streamlit app · trained on real ASTRAM / BTP event data.", 14,
        GREY, False)]])

# ---- Slide 4: Forecast figure ---------------------------------------------
s = slide()
header(s, "FORECAST", "Why a range beats a single number")
pic(s, os.path.join(FIG, "fig1_forecast_intervals.png"), 1.6, 2.2, 10.1)

# ---- Slide 5: Model + CQR -------------------------------------------------
s = slide()
header(s, "THE MODEL", "Calibrated, leakage-safe, honest")
bullets(s, [
    "LightGBM on log-resolution + quantile models (P10/P50/P90).",
    "Dedicated models for high-variance causes (potholes, construction).",
    "Conformalized Quantile Regression calibrates the intervals.",
    "Strict temporal validation — train on earlier events, test on later.",
], y=2.5, dy=0.82, size=15)
pic(s, os.path.join(FIG, "fig3_cqr_coverage.png"), 8.0, 2.6, 4.6)

# ---- Slide 6: Dispatch heatmap --------------------------------------------
s = slide()
header(s, "DISPATCH", "Where & when to pre-position crews")
pic(s, os.path.join(FIG, "fig2_demand_heatmap.png"), 1.4, 2.15, 10.5)

# ---- Slide 7: Performance -------------------------------------------------
s = slide()
header(s, "PERFORMANCE", "Honest numbers (temporal hold-out)")
cards = [(f"{M['point_medae_h']:.1f} h", "Typical (median) error"),
         ("~0.7 h", "Vehicle breakdown P50 (60% of all events)"),
         (f"{M['p10_p90_coverage']:.0%}", "P10–P90 coverage (60% → after CQR)")]
for i, (big, small) in enumerate(cards):
    x = 1.0 + i * 3.95
    rect(s, x, 2.6, 3.5, 1.9, CARD)
    text(s, x, 2.85, 3.5, 0.9, [[(big, 36, BLUE, True)]], align=PP_ALIGN.CENTER)
    text(s, x + 0.2, 3.75, 3.1, 0.7, [[(small, 13, GREY, False)]],
         align=PP_ALIGN.CENTER)
text(s, 1.0, 5.0, 11.3, 1.3,
     [[("Mean error looks large only because rare multi-day infrastructure "
        "incidents dominate the average — which is exactly why we report "
        "intervals. Surfaced in the app's Model Card, not hidden.", 15, DARK,
        False)]])

# ---- Slide 8: Demo + close ------------------------------------------------
s = slide()
header(s, "DEMO & IMPACT", "From incident to action in one screen")
bullets(s, [
    "Breakdown, Mysore Rd, 6 AM → ~0.7 h, HIGH, 2 towing units.",
    "Waterlogging + closure → ~24 h (P90 ~170 h), CRITICAL, pumps + diversion.",
    "Dispatch tab → demand heatmap + downloadable staging plan.",
    "Hotspots tab → preventive-maintenance targets.",
], y=2.4, dy=0.78, size=16)
text(s, 1.0, 5.8, 11.3, 0.6,
     [[("Built on real BTP data · deployable today · scales to live ASTRAM feeds.",
        16, NAVY, True)]])
text(s, 1.0, 6.5, 11.3, 0.6, [[("Thank you.", 18, BLUE, True)]])

prs.save(OUT)
print(f"Wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")
